"""
Office 文件工具（读写 docx/xlsx/pptx/pdf）
依赖（按需安装，缺少时优雅降级）：
  pip install python-docx openpyxl python-pptx reportlab pdfplumber Pillow
"""

import os
import io
import json
from pathlib import Path
from typing import Dict, Any, List, Optional


def _get_desktop() -> Path:
    """Get user Desktop folder (cross-platform)"""
    import sys, subprocess
    p = Path.home() / "Desktop"
    if p.exists():
        return p
    if sys.platform == "linux":
        try:
            result = subprocess.run(
                ["xdg-user-dir", "DESKTOP"],
                capture_output=True, text=True, timeout=3
            )
            if result.returncode == 0 and result.stdout.strip():
                alt = Path(result.stdout.strip())
                if alt.exists():
                    return alt
        except Exception:
            pass
    return Path.home()


# ══════════════════════════════════════════
# 读取工具
# ══════════════════════════════════════════

def read_docx(path: str) -> Dict:
    """读取 Word 文档（.docx）"""
    try:
        from docx import Document
        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = []
        for t in doc.tables:
            rows = [[cell.text for cell in row.cells] for row in t.rows]
            tables.append(rows)
        return {
            "ok": True, "type": "docx",
            "paragraphs": paragraphs,
            "text": "\n".join(paragraphs),
            "tables": tables,
            "table_count": len(tables),
            "paragraph_count": len(paragraphs)
        }
    except ImportError:
        return {"ok": False, "error": "需要安装 python-docx：pip install python-docx"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def read_xlsx(path: str) -> Dict:
    """读取 Excel 文档（.xlsx/.xls）"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheets = {}
        for name in wb.sheetnames:
            ws = wb[name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(c) if c is not None else "" for c in row]
                if any(row_data):
                    rows.append(row_data)
            sheets[name] = rows
        return {
            "ok": True, "type": "xlsx",
            "sheets": sheets,
            "sheet_names": wb.sheetnames,
            "text": "\n\n".join(
                f"[Sheet: {n}]\n" + "\n".join("\t".join(r) for r in rows)
                for n, rows in sheets.items()
            )
        }
    except ImportError:
        return {"ok": False, "error": "需要安装 openpyxl：pip install openpyxl"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def read_pptx(path: str) -> Dict:
    """读取 PowerPoint 文档（.pptx）"""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            slides.append({"slide": i, "texts": texts})
        all_text = "\n\n".join(
            f"[Slide {s['slide']}]\n" + "\n".join(s['texts'])
            for s in slides if s['texts']
        )
        return {
            "ok": True, "type": "pptx",
            "slides": slides,
            "slide_count": len(prs.slides),
            "text": all_text
        }
    except ImportError:
        return {"ok": False, "error": "需要安装 python-pptx：pip install python-pptx"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def read_pdf(path: str) -> Dict:
    """读取 PDF 文档"""
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            pages = []
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                tables = page.extract_tables() or []
                pages.append({"page": i, "text": text, "tables": tables})
        all_text = "\n\n".join(
            f"[Page {p['page']}]\n{p['text']}" for p in pages if p['text']
        )
        return {
            "ok": True, "type": "pdf",
            "pages": pages,
            "page_count": len(pages),
            "text": all_text
        }
    except ImportError:
        # 备用：用 PyPDF2
        try:
            import PyPDF2
            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                texts = [page.extract_text() or "" for page in reader.pages]
            return {
                "ok": True, "type": "pdf",
                "page_count": len(texts),
                "text": "\n\n".join(f"[Page {i+1}]\n{t}" for i, t in enumerate(texts))
            }
        except ImportError:
            return {"ok": False, "error": "需要安装 pdfplumber：pip install pdfplumber"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def read_office_file(path: str) -> Dict:
    """
    自动检测文件类型并读取
    支持：.docx .doc .xlsx .xls .pptx .ppt .pdf .txt .md .csv
    """
    path = os.path.expanduser(path)
    if not os.path.exists(path):
        return {"ok": False, "error": f"文件不存在：{path}"}

    ext = Path(path).suffix.lower()
    if ext in (".docx", ".doc"):
        return read_docx(path)
    elif ext in (".xlsx", ".xls", ".csv"):
        if ext == ".csv":
            try:
                with open(path, encoding="utf-8", errors="replace") as f:
                    content = f.read()
                return {"ok": True, "type": "csv", "text": content[:30000]}
            except Exception as e:
                return {"ok": False, "error": str(e)}
        return read_xlsx(path)
    elif ext in (".pptx", ".ppt"):
        return read_pptx(path)
    elif ext == ".pdf":
        return read_pdf(path)
    elif ext in (".txt", ".md", ".py", ".js", ".json", ".html", ".css"):
        try:
            with open(path, encoding="utf-8", errors="replace") as f:
                content = f.read()
            return {"ok": True, "type": ext.lstrip("."), "text": content[:50000]}
        except Exception as e:
            return {"ok": False, "error": str(e)}
    else:
        return {"ok": False, "error": f"不支持的文件类型：{ext}"}


# ══════════════════════════════════════════
# 创建工具
# ══════════════════════════════════════════

def create_docx(path: str, content: str, title: str = "") -> Dict:
    """
    创建 Word 文档
    content 支持简单 Markdown：# 标题，**粗体**，- 列表，| 表格
    """
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # 设置默认字体
        style = doc.styles["Normal"]
        style.font.name = "Arial"
        style.font.size = Pt(11)

        if title:
            h = doc.add_heading(title, level=0)
            h.alignment = WD_ALIGN_PARAGRAPH.CENTER

        for line in content.split("\n"):
            stripped = line.strip()
            if not stripped:
                doc.add_paragraph("")
                continue
            # 标题
            if stripped.startswith("### "):
                doc.add_heading(stripped[4:], level=3)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:], level=2)
            elif stripped.startswith("# "):
                doc.add_heading(stripped[2:], level=1)
            # 列表
            elif stripped.startswith("- ") or stripped.startswith("* "):
                doc.add_paragraph(stripped[2:], style="List Bullet")
            elif stripped[:2].isdigit() and stripped[2] in (". ", ") "):
                doc.add_paragraph(stripped[3:], style="List Number")
            # 表格行（简单处理）
            elif stripped.startswith("|") and stripped.endswith("|"):
                cells = [c.strip() for c in stripped.strip("|").split("|")]
                # 跳过分隔行
                if all(set(c) <= set("-: ") for c in cells):
                    continue
                # 获取或创建表格
                if not doc.tables or not hasattr(doc.tables[-1], "_cells_added"):
                    t = doc.add_table(rows=0, cols=len(cells))
                    t.style = "Table Grid"
                    t._cells_added = True
                else:
                    t = doc.tables[-1]
                row = t.add_row()
                for i, cell_text in enumerate(cells):
                    if i < len(row.cells):
                        row.cells[i].text = cell_text
            else:
                # 处理粗体 **text**
                para = doc.add_paragraph()
                import re
                parts = re.split(r'\*\*(.+?)\*\*', stripped)
                for j, part in enumerate(parts):
                    run = para.add_run(part)
                    if j % 2 == 1:
                        run.bold = True

        # 确保目录存在
        path = _resolve_output_path(path)
        doc.save(path)
        return {"ok": True, "path": path, "type": "docx"}
    except ImportError:
        return {"ok": False, "error": "需要安装 python-docx：pip install python-docx"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def create_xlsx(path: str, data: Any, sheet_name: str = "Sheet1") -> Dict:
    """
    创建 Excel 文档
    data: List[List] 或 dict {"sheet_name": [[row], ...]}
    """
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment

        wb = openpyxl.Workbook()
        wb.remove(wb.active)

        if isinstance(data, list):
            data = {sheet_name: data}

        for sname, rows in data.items():
            ws = wb.create_sheet(sname)
            for i, row in enumerate(rows, 1):
                for j, val in enumerate(row, 1):
                    cell = ws.cell(row=i, column=j, value=val)
                    # 第一行加粗（标题行）
                    if i == 1:
                        cell.font = Font(bold=True)
                        cell.fill = PatternFill("solid", fgColor="D5E8F0")
                        cell.alignment = Alignment(horizontal="center")
            # 自动列宽
            for col in ws.columns:
                max_len = max((len(str(c.value or "")) for c in col), default=0)
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

        path = _resolve_output_path(path)
        wb.save(path)
        return {"ok": True, "path": path, "type": "xlsx"}
    except ImportError:
        return {"ok": False, "error": "需要安装 openpyxl：pip install openpyxl"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def create_pptx(path: str, slides_data: List[Dict]) -> Dict:
    """
    创建 PowerPoint 文档
    slides_data: [{"title": "...", "content": "...", "bullets": [...]}]
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        # 16:9
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        for slide_data in slides_data:
            title_text   = slide_data.get("title", "")
            content_text = slide_data.get("content", "")
            bullets      = slide_data.get("bullets", [])

            if title_text and (content_text or bullets):
                layout = prs.slide_layouts[1]  # Title and Content
            elif title_text:
                layout = prs.slide_layouts[0]  # Title Slide
            else:
                layout = prs.slide_layouts[6]  # Blank

            slide = prs.slides.add_slide(layout)

            # 标题
            if slide.shapes.title and title_text:
                slide.shapes.title.text = title_text
                tf = slide.shapes.title.text_frame
                tf.paragraphs[0].runs[0].font.bold = True
                tf.paragraphs[0].runs[0].font.size = Pt(32)

            # 内容
            body_texts = []
            if content_text:
                body_texts.append(content_text)
            body_texts.extend(bullets)

            if body_texts and len(slide.placeholders) > 1:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for i, bt in enumerate(body_texts):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = bt
                    p.font.size = Pt(18)

        path = _resolve_output_path(path)
        prs.save(path)
        return {"ok": True, "path": path, "type": "pptx",
                "slide_count": len(slides_data)}
    except ImportError:
        return {"ok": False, "error": "需要安装 python-pptx：pip install python-pptx"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def create_pdf(path: str, content: str, title: str = "") -> Dict:
    """
    创建 PDF 文档（支持中文）
    """
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        import textwrap

        path = _resolve_output_path(path)
        c = canvas.Canvas(path, pagesize=A4)
        width, height = A4

        # 尝试注册中文字体
        font_name = "Helvetica"
        for font_path in [
            "/System/Library/Fonts/PingFang.ttc",
            "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
            "C:/Windows/Fonts/msyh.ttc",
            "C:/Windows/Fonts/simhei.ttf",
        ]:
            if os.path.exists(font_path):
                try:
                    pdfmetrics.registerFont(TTFont("CJK", font_path))
                    font_name = "CJK"
                    break
                except Exception:
                    pass

        margin = 20 * mm
        y = height - margin
        line_height = 6 * mm

        if title:
            c.setFont(font_name, 18)
            c.drawString(margin, y, title)
            y -= line_height * 2

        c.setFont(font_name, 11)
        for line in content.split("\n"):
            if y < margin + line_height:
                c.showPage()
                y = height - margin
                c.setFont(font_name, 11)
            # 处理标题
            if line.startswith("# "):
                c.setFont(font_name, 16)
                c.drawString(margin, y, line[2:])
                c.setFont(font_name, 11)
            elif line.startswith("## "):
                c.setFont(font_name, 14)
                c.drawString(margin, y, line[3:])
                c.setFont(font_name, 11)
            else:
                # 自动换行（英文按字符，中文也按字符）
                wrapped = textwrap.wrap(line, width=80) if line else [""]
                for wline in (wrapped or [""]):
                    c.drawString(margin, y, wline)
                    y -= line_height
                    if y < margin + line_height:
                        c.showPage()
                        y = height - margin
                        c.setFont(font_name, 11)
                continue
            y -= line_height

        c.save()
        return {"ok": True, "path": path, "type": "pdf"}
    except ImportError:
        return {"ok": False, "error": "需要安装 reportlab：pip install reportlab"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


# ══════════════════════════════════════════
# 图片识别（Vision）
# ══════════════════════════════════════════

def analyze_image(image_path: str, question: str = "",
                  api_key: str = "", provider: str = "openai") -> Dict:
    """
    用 LLM Vision API 分析图片
    支持 OpenAI GPT-4V / Anthropic Claude / Google Gemini
    """
    import base64

    if not os.path.exists(image_path):
        return {"ok": False, "error": f"图片不存在：{image_path}"}

    with open(image_path, "rb") as f:
        img_b64 = base64.b64encode(f.read()).decode()

    ext = Path(image_path).suffix.lower().lstrip(".")
    media_type = {"jpg": "image/jpeg", "jpeg": "image/jpeg",
                  "png": "image/png", "gif": "image/gif",
                  "webp": "image/webp"}.get(ext, "image/jpeg")

    prompt = question or "请描述这张图片的内容，包括主要对象、场景、文字等关键信息。"

    try:
        if provider in ("openai", "deepseek"):
            # OpenAI Vision 格式
            import urllib.request
            base_url = ("https://api.deepseek.com/v1" if provider == "deepseek"
                        else "https://api.openai.com/v1")
            model = "deepseek-chat" if provider == "deepseek" else "gpt-4o-mini"
            payload = json.dumps({
                "model": model,
                "messages": [{
                    "role": "user",
                    "content": [
                        {"type": "image_url",
                         "image_url": {"url": f"data:{media_type};base64,{img_b64}"}},
                        {"type": "text", "text": prompt}
                    ]
                }],
                "max_tokens": 1000
            }).encode()
            req = urllib.request.Request(
                f"{base_url}/chat/completions", data=payload,
                headers={"Content-Type": "application/json",
                         "Authorization": f"Bearer {api_key}"}
            )
            with urllib.request.urlopen(req, timeout=30) as r:
                data = json.loads(r.read())
            return {"ok": True, "description": data["choices"][0]["message"]["content"]}

        elif provider == "claude":
            import urllib.request
            payload = json.dumps({
                "model": "claude-3-5-haiku-20241022",
                "max_tokens": 1000,
                "messages": [{
                    "role": "user",
                    "content": [
                        {"type": "image",
                         "source": {"type": "base64", "media_type": media_type,
                                    "data": img_b64}},
                        {"type": "text", "text": prompt}
                    ]
                }]
            }).encode()
            req = urllib.request.Request(
                "https://api.anthropic.com/v1/messages", data=payload,
                headers={"Content-Type": "application/json",
                         "x-api-key": api_key,
                         "anthropic-version": "2023-06-01"}
            )
            with urllib.request.urlopen(req, timeout=30) as r:
                data = json.loads(r.read())
            return {"ok": True, "description": data["content"][0]["text"]}

        elif provider == "gemini":
            import urllib.request
            payload = json.dumps({
                "contents": [{
                    "parts": [
                        {"inline_data": {"mime_type": media_type, "data": img_b64}},
                        {"text": prompt}
                    ]
                }]
            }).encode()
            url = (f"https://generativelanguage.googleapis.com/v1beta/models/"
                   f"gemini-1.5-flash:generateContent?key={api_key}")
            req = urllib.request.Request(url, data=payload,
                                         headers={"Content-Type": "application/json"})
            with urllib.request.urlopen(req, timeout=30) as r:
                data = json.loads(r.read())
            return {"ok": True,
                    "description": data["candidates"][0]["content"]["parts"][0]["text"]}

    except Exception as e:
        return {"ok": False, "error": str(e)}

    return {"ok": False, "error": f"不支持的 provider: {provider}"}


# ══════════════════════════════════════════
# 工具函数
# ══════════════════════════════════════════

def _resolve_output_path(path: str) -> str:
    """如果只是文件名，保存到桌面"""
    path = path.strip()
    if not any(c in path for c in ["/", "\\", ":"]):
        desktop = _get_desktop()
        path = str(desktop / path)
    path = os.path.expanduser(path)
    os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
    return path


def check_dependencies() -> Dict[str, bool]:
    """检查各 Office 库是否安装"""
    deps = {}
    for pkg, import_name in [
        ("python-docx", "docx"),
        ("openpyxl", "openpyxl"),
        ("python-pptx", "pptx"),
        ("pdfplumber", "pdfplumber"),
        ("reportlab", "reportlab"),
        ("Pillow", "PIL"),
    ]:
        try:
            __import__(import_name)
            deps[pkg] = True
        except ImportError:
            deps[pkg] = False
    return deps
